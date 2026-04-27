"""Build the Open University Google Trends deck. Brand-compliant SumItUp PPTX."""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image


# ---------- Brand constants ----------
class BRAND:
    NAVY = RGBColor(0x28, 0x22, 0x6C)
    PURPLE = RGBColor(0x59, 0x40, 0xC0)
    TEAL = RGBColor(0x04, 0x93, 0xA3)
    GOLD = RGBColor(0xFF, 0xC6, 0x55)
    CREAM = RGBColor(0xF9, 0xF4, 0xEE)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    FONT = "Heebo"
    SLIDE_W = Inches(13.33)
    SLIDE_H = Inches(7.5)


def set_rtl(tf, right_align=True):
    for para in tf.paragraphs:
        pPr = para._p.get_or_add_pPr()
        pPr.set("rtl", "1")
        if right_align:
            para.alignment = PP_ALIGN.RIGHT


def apply_font(run, size=15, bold=False, color=BRAND.NAVY, name=BRAND.FONT):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    cs = rPr.find(qn("a:cs"))
    if cs is None:
        cs = etree.SubElement(rPr, qn("a:cs"))
    cs.set("typeface", name)


def fill_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def compute_fit(img_path, max_w=12.33, max_h=4.95, slide_w=13.33):
    with Image.open(img_path) as im:
        iw, ih = im.size
    img_r = iw / ih
    box_r = max_w / max_h
    if img_r > box_r:
        w = max_w
        h = w / img_r
    else:
        h = max_h
        w = h * img_r
    x = (slide_w - w) / 2
    return x, w, h


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = BRAND.SLIDE_W
        self.prs.slide_height = BRAND.SLIDE_H
        self.page = 0

    # ---- footer ----
    def _footer(self, slide):
        line = slide.shapes.add_connector(1, Inches(0.5), Inches(7.05), Inches(12.83), Inches(7.05))
        line.line.color.rgb = BRAND.NAVY
        line.line.width = Pt(1)

        tb = slide.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(1.5), Inches(0.35))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p._p.get_or_add_pPr().set("rtl", "0")
        r = p.add_run()
        r.text = "sumitup."
        apply_font(r, size=18, bold=True)

        tb = slide.shapes.add_textbox(Inches(12.0), Inches(7.1), Inches(1.0), Inches(0.25))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        p._p.get_or_add_pPr().set("rtl", "0")
        r = p.add_run()
        r.text = str(self.page)
        apply_font(r, size=11)

    def _source(self, slide, text):
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(6.75), Inches(8.0), Inches(0.22))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p._p.get_or_add_pPr().set("rtl", "0")
        r = p.add_run()
        r.text = text
        apply_font(r, size=10)

    # ---- cover ----
    def cover(self, title, subtitle):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        fill_bg(s, BRAND.NAVY)

        tb = s.shapes.add_textbox(0, Inches(2.4), BRAND.SLIDE_W, Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = title
        apply_font(r, size=54, bold=True, color=BRAND.WHITE)

        tb = s.shapes.add_textbox(0, Inches(4.4), BRAND.SLIDE_W, Inches(0.8))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = subtitle
        apply_font(r, size=28, bold=False, color=BRAND.WHITE)

        return s

    # ---- section divider ----
    def section(self, title):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        fill_bg(s, BRAND.NAVY)
        tb = s.shapes.add_textbox(0, Inches(3.2), BRAND.SLIDE_W, Inches(1.3))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = title
        apply_font(r, size=44, bold=True, color=BRAND.WHITE)
        return s

    # ---- content with title + bullets ----
    def content(self, title, bullets, source=None, insight=None):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        fill_bg(s, BRAND.CREAM)
        self.page += 1

        # Title
        tb = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(0.7))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title
        apply_font(r, size=26, bold=True)

        y_start = 1.5
        # Optional insight headline (teal, bigger)
        if insight:
            tb = s.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.33), Inches(0.8))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = insight
            apply_font(r, size=20, bold=True, color=BRAND.TEAL)
            y_start = 2.4

        # Body bullets
        tb = s.shapes.add_textbox(Inches(0.5), Inches(y_start), Inches(12.33), Inches(6.5 - y_start))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(10)
            r = p.add_run()
            r.text = "• " + line
            apply_font(r, size=15)

        if source:
            self._source(s, source)
        self._footer(s)
        return s

    # ---- content with title + chart image ----
    def chart_slide(self, title, chart_path, caption=None, source=None, insight=None):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        fill_bg(s, BRAND.CREAM)
        self.page += 1

        tb = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(0.7))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title
        apply_font(r, size=26, bold=True)

        y_img_start = 1.5
        if insight:
            tb = s.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.33), Inches(0.7))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = insight
            apply_font(r, size=18, bold=True, color=BRAND.TEAL)
            y_img_start = 2.2

        # Chart
        x, w, h = compute_fit(chart_path, max_w=11.5, max_h=6.3 - y_img_start)
        pic = s.shapes.add_picture(chart_path, Inches(x), Inches(y_img_start), Inches(w), Inches(h))

        # Teal frame
        frame = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x - 0.08), Inches(y_img_start - 0.08),
            Inches(w + 0.16), Inches(h + 0.16),
        )
        frame.fill.background()
        frame.line.color.rgb = BRAND.TEAL
        frame.line.width = Pt(5)
        frame.rotation = -1.5

        if caption:
            cap_y = y_img_start + h + 0.15
            if cap_y < 6.6:
                tb = s.shapes.add_textbox(Inches(0.5), Inches(cap_y), Inches(12.33), Inches(0.3))
                p = tb.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                r = p.add_run()
                r.text = caption
                apply_font(r, size=12, color=BRAND.NAVY)

        if source:
            self._source(s, source)
        self._footer(s)
        return s

    # ---- big stat ----
    def big_stat(self, stat, label, color=None, source=None):
        color = color or BRAND.TEAL
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        fill_bg(s, BRAND.CREAM)
        self.page += 1

        tb = s.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(12.33), Inches(2.4))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = stat
        apply_font(r, size=120, bold=True, color=color)

        tb = s.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(12.33), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        apply_font(r, size=22, color=BRAND.NAVY)

        if source:
            self._source(s, source)
        self._footer(s)
        return s

    # ---- Hebrew pull-quote (centered large RTL) ----
    def hebrew_quote(self, hebrew_text, english_caption, color=None, source=None):
        color = color or BRAND.PURPLE
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        fill_bg(s, BRAND.CREAM)
        self.page += 1

        tb = s.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(12.33), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p._p.get_or_add_pPr().set("rtl", "1")
        r = p.add_run()
        r.text = '"' + hebrew_text + '"'
        apply_font(r, size=64, bold=True, color=color)

        tb = s.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(12.33), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = english_caption
        apply_font(r, size=20, color=BRAND.NAVY)

        if source:
            self._source(s, source)
        self._footer(s)
        return s

    # ---- closing ----
    def closing(self, text="Thank you."):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        fill_bg(s, BRAND.NAVY)
        tb = s.shapes.add_textbox(0, Inches(3.2), BRAND.SLIDE_W, Inches(1.3))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        apply_font(r, size=44, bold=True, color=BRAND.WHITE)
        return s

    def save(self, path):
        self.prs.save(path)


# ================================================================
# DECK CONTENT
# ================================================================
def build(out_path: Path, run_dir: Path):
    charts = run_dir / "charts"
    d = Deck()

    # 1. Cover
    d.cover(
        title="Open University × Higher Education",
        subtitle="Google Trends Read — April 2026",
    )

    # 2. The question
    d.content(
        title="The question we tried to answer",
        insight="Is interest in higher education growing — and is the motivation shifting?",
        bullets=[
            "Traditional logic: study → get a credential (תעודה) → get a job.",
            "The client's hypothesis: this reasoning is eroding among young Israelis.",
            "We read 5 years of Google Trends data across brand, category, and adjacent queries to test it.",
            "Scope: Israel, Hebrew, 5-year timeframe, ~35 query bundles.",
        ],
    )

    # 3. Executive summary
    d.content(
        title="Executive summary",
        bullets=[
            "Interest in higher education isn't collapsing — but Open U's share of it is.",
            "The credential→job logic is fraying at the edges, not being overthrown. The real shift is toward cost anxiety and gate-skipping, not degree-rejection.",
            "The category is fragmenting into logistics and specialization — people ask how many credit points, what it costs, which campus — not whether a degree is worth it.",
            "Open U's strategic opening is psychometric-free, cost-transparent, indecision-friendly positioning — not an 'alternative to a degree' story.",
        ],
    )

    # 4. Section: Brand Position
    d.section("Brand Position")

    # 5. Brand vs competitors chart
    d.chart_slide(
        title="Open U is not in the competitive conversation",
        insight="TAU dominates 75.5% SoV. Open U is below the head-to-head threshold.",
        chart_path=str(charts / "brand_vs_competitors.png"),
        source="Google Trends, Israel, 5-year, Hebrew-only queries.",
    )

    # 6. Big stat: -7.6% YoY
    d.big_stat(
        stat="-7.6%",
        label='YoY decline on "האוניברסיטה הפתוחה" (Open University) brand search, last 52 weeks.',
        color=BRAND.PURPLE,
        source="Google Trends, last 52 weeks vs prior 52 weeks.",
    )

    # 7. Brand terms trend
    d.chart_slide(
        title="Open U brand interest — a slow-grade erosion",
        insight="Primary brand term baseline: 58.2 → 53.8 over the last 52w. Long-term slope: -3.74/year.",
        chart_path=str(charts / "brand_terms.png"),
        source="Google Trends, Hebrew brand terms for Open University.",
    )

    # 8. Seasonality
    d.chart_slide(
        title="The July enrollment pulse still works",
        insight="Seasonality strength 0.41. July index 1.2x — timing the annual enrollment window is right.",
        chart_path=str(charts / "brand_seasonality.png"),
        source="Google Trends, Open University brand term, monthly seasonality.",
    )

    # 9. Section: Category Context
    d.section("Category Context")

    # 10. Where attention is flowing
    d.content(
        title="The HE category is a mild tailwind — with a compositional shift",
        insight="Energy is moving from 'should I study?' to 'what are the exact rules, the exact field, the exact campus.'",
        bullets=[
            '"תחילת לימודים אקדמיים 2023" (academic year start) rising +277,550% — the annual enrollment pulse is alive.',
            '"נקודות זיכוי תואר ראשון" (credit points), "כמה שנים תואר ראשון" (degree length) — demand is moving to program mechanics.',
            '"תואר ראשון בפיזיותרפיה" (BA in physiotherapy) +50 — field specialization is growing.',
            '"תל חי אוניברסיטה" +250 — non-flagship exploration. Students are looking beyond the obvious brands.',
        ],
        source="Google Trends rising related queries for category terms.",
    )

    # 11. Adjacent paths chart
    d.chart_slide(
        title="The 'alternative path' vocabulary exists — but volume is thin",
        insight='100% of last-52w SoV in this cluster is "לימודים אונליין" (online studies). Zero for anti-degree framings.',
        chart_path=str(charts / "adjacent_1.png"),
        source="Google Trends: online studies, hi-tech without degree, profession with a future.",
    )

    # 12. Degree vs bootcamp
    d.chart_slide(
        title='"תואר ראשון" (BA) vs "בוטקמפ" (bootcamp) — no contest',
        insight="The bootcamp narrative is loud in media. In search behavior it barely moves the needle.",
        chart_path=str(charts / "compare_1.png"),
        source="Google Trends head-to-head, Israel, 5-year.",
    )

    # 13. Section: The Audience Story
    d.section("The Audience Story")

    # 14. Hebrew quote — indecision
    d.hebrew_quote(
        hebrew_text="לא יודע/ת מה ללמוד",
        english_caption="Young Israelis are literally googling their own indecision.",
        color=BRAND.PURPLE,
        source="Top rising queries in category, Google Trends.",
    )

    # 15. The "why" question answered
    d.content(
        title='The "why" is not what the client thought',
        insight="The credential→job logic isn't being overthrown — it's being stress-tested at the cost and entry-gate seams.",
        bullets=[
            '"כמה עולה תואר ראשון" (how much does a BA cost) and "נקודת זיכוי" (tax credit) — top rising. The ROI frame became a cost frame.',
            '"לימודים אקדמיים ללא פסיכומטרי" (academic studies without psychometric) — top query. People want the credential without the traditional gate.',
            '"דרושים תואר ראשון" (jobs requiring a BA) still ranks top. The job-market gate is still perceived as real.',
            "Almost no queries for bootcamps vs degrees, no reviews, no 'best' rankings. Decisions are institution-first, not path-first.",
            "Two segments: post-army freshmen (Sep spike + deferral queries) and women in retraining. Neither rejects credentials.",
        ],
    )

    # 16. Section: Strategic Implications
    d.section("Strategic Implications")

    # 17. Implications
    d.content(
        title="Five moves for Open University",
        bullets=[
            '1. Own "credential without the gate" — "לימודים אקדמיים ללא פסיכומטרי" is a top query and Open U\'s structural advantage. Build a content pillar + SEO + paid search around it.',
            '2. Build the missing cost calculator. "כמה עולה תואר ראשון" and "נקודת זיכוי" have no authoritative owner. Turn financial anxiety into a lead-gen asset.',
            '3. Answer the indecision query directly. "לא יודע/ת מה ללמוד" is a gap nobody is filling. A guided "what should I study" diagnostic is brand-building and demand-capture in one move.',
            '4. Fix the portal now. "אזור אישי" +700% and "כניסה" +110% aren\'t brand signals — they\'re UX complaints from enrolled students. Unsexy but high-ROI.',
            '5. Expand "אקדמיה בתיכון" and the deferral cohort. High-school funnel +120% for Open U. Reservists need flexible scheduling — Open U\'s literal product.',
        ],
    )

    # 18. Data caveats
    d.content(
        title="Data caveats",
        bullets=[
            "IOT fetches failed for several chunks due to Google Trends rate limits; some 'alternative path' series are thin by volume, not necessarily by underlying interest.",
            "All numbers are relative search interest (0–100). We can compare trajectories and share, not absolute query counts or enrollments.",
            "Niche queries near the 0–1 threshold are directionally honest but statistically noisy — treat single-week spikes as weak signal.",
            "Open U is absent from the competitor head-to-head because Google Trends returned it below threshold — itself a finding, but worth triangulating against paid search impression share before acting on it.",
        ],
    )

    # 19. Closing
    d.closing("Let's discuss.")

    d.save(str(out_path))
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    run_dir = Path(sys.argv[1]) if len(sys.argv) > 1 else Path("data/open_university_20260422_191112")
    out = run_dir / "deck.pptx"
    build(out, run_dir)
