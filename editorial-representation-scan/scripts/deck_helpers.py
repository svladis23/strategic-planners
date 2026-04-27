"""
SumItUp brand-compliant slide template functions for the Editorial Representation Scan.
Called by scripts/render_deck.py after reading deck_plan.yaml.

Brand specs (from sumitup-powerpoint skill):
- Slide size: 13.33 × 7.5 in (16:9)
- Palette: Navy #28226C / Purple #5940C0 / Teal #0493A3 / Gold #FFC655 / Cream #F9F4EE / White / Black
- Font: Heebo (Bold / Medium / Regular / Light)
- Cover: navy bg. Content: cream bg. Dividers: navy.
- Footer on every content slide: "sumitup." left + separator line + page num right.
- Hebrew: RTL, right-align. English: LTR, left-align.
- No drop shadows, no gradients, no filled header/footer bars beyond the thin separator.

Each template function takes (prs, page_counter, params: dict) and appends ONE slide.
The renderer increments page_counter and passes it to add_footer() where appropriate.

To add a new template:
1. Add a function below with signature (prs, page_counter, params) -> int (new page_counter).
2. Register it in TEMPLATES dict at the bottom.
3. Document it in docs/slide_templates.md.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree


# ============================================================
# BRAND CONSTANTS
# ============================================================
class B:
    NAVY   = RGBColor(0x28, 0x22, 0x6C)
    PURPLE = RGBColor(0x59, 0x40, 0xC0)
    TEAL   = RGBColor(0x04, 0x93, 0xA3)
    GOLD   = RGBColor(0xFF, 0xC6, 0x55)
    CREAM  = RGBColor(0xF9, 0xF4, 0xEE)
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    BLACK  = RGBColor(0x00, 0x00, 0x00)
    GRAY   = RGBColor(0xAA, 0xAA, 0xAA)
    FONT   = "Heebo"


TAG_COLORS = {
    "CONVERGENT": B.NAVY,
    "DIVERGENT":  B.PURPLE,
    "RISING":     B.TEAL,
    "COUNTER":    B.GOLD,
    "STABLE":     B.NAVY,
    "DECLINING":  B.GRAY,
}


# ============================================================
# LOW-LEVEL HELPERS
# ============================================================

def _apply_font(run, size=15, bold=False, color=B.NAVY, italic=False, name=B.FONT):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:cs', 'a:ea'):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set('typeface', name)


def _fill_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, x, y, w, h, text, size=15, bold=False, color=B.NAVY,
              align=PP_ALIGN.LEFT, rtl=False, italic=False, anchor=None, wrap=True):
    """Add a textbox with single-format text. `text` can be str or list of str (one per line)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    if anchor is not None:
        tf.vertical_anchor = anchor
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if rtl:
            pPr = p._p.get_or_add_pPr()
            pPr.set('rtl', '1')
        r = p.add_run()
        r.text = line
        _apply_font(r, size=size, bold=bold, color=color, italic=italic)
    return tb


def _add_footer(slide, page_num):
    """sumitup. logo + separator line + page number."""
    # Separator line
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(7.05), Inches(12.83), Inches(7.05))
    line.line.color.rgb = B.NAVY
    line.line.width = Pt(1)
    # Logo (LTR)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(1.5), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    pPr = p._p.get_or_add_pPr(); pPr.set('rtl', '0'); pPr.set('algn', 'l')
    r = p.add_run(); r.text = "sumitup."
    _apply_font(r, size=18, bold=True, color=B.NAVY)
    # Page number (LTR, right)
    tb = slide.shapes.add_textbox(Inches(12.0), Inches(7.1), Inches(1.0), Inches(0.25))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    pPr = p._p.get_or_add_pPr(); pPr.set('rtl', '0'); pPr.set('algn', 'r')
    r = p.add_run(); r.text = str(page_num)
    _apply_font(r, size=11, color=B.NAVY)


def _add_source_citation(slide, text):
    """Source line right above the footer separator."""
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(6.75), Inches(12.33), Inches(0.22))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    pPr = p._p.get_or_add_pPr(); pPr.set('rtl', '0'); pPr.set('algn', 'l')
    r = p.add_run(); r.text = text
    _apply_font(r, size=9, color=B.NAVY, italic=True)


def _add_title(slide, text, size=24):
    _add_text(slide, 0.5, 0.5, 12.33, 0.75, text,
              size=size, bold=True, color=B.NAVY, align=PP_ALIGN.LEFT)


def _add_subhead(slide, text, y=1.25, size=12, color=None, italic=True):
    _add_text(slide, 0.5, y, 12.33, 0.4, text,
              size=size, italic=italic, color=color or B.TEAL,
              align=PP_ALIGN.LEFT, bold=True)


def _hebrew_quote_block(slide, x, y, w, hebrew, gloss, attribution, h_heb=0.9):
    """Standard 3-part quote block: Hebrew RTL + English italic gloss + attribution."""
    # Hebrew (RTL, navy bold)
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h_heb))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    pPr = p._p.get_or_add_pPr(); pPr.set('rtl', '1')
    r = p.add_run(); r.text = '"' + hebrew + '"'
    _apply_font(r, size=13, bold=True, color=B.NAVY)
    # Gloss (LTR italic)
    tb = slide.shapes.add_textbox(Inches(x), Inches(y + h_heb), Inches(w), Inches(0.5))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = '"' + gloss + '"'
    _apply_font(r, size=11, italic=True, color=B.NAVY)
    # Attribution
    tb = slide.shapes.add_textbox(Inches(x), Inches(y + h_heb + 0.45), Inches(w), Inches(0.3))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "— " + attribution
    _apply_font(r, size=10, bold=True, color=B.TEAL)


def _new_content_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide, B.CREAM)
    return slide


def _new_navy_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide, B.NAVY)
    return slide


# ============================================================
# TEMPLATE FUNCTIONS
# ============================================================
# Each template takes (prs, page_num, params) and returns the next page_num.
# page_num==0 means "don't add footer" (covers and section dividers).
# ============================================================

def tpl_cover(prs, page_num, p):
    s = _new_navy_slide(prs)
    _add_text(s, 0, 1.5, 13.33, 0.4,
              "EDITORIAL REPRESENTATION SCAN",
              size=16, color=B.TEAL, align=PP_ALIGN.CENTER, bold=True)
    _add_text(s, 0.5, 2.3, 12.33, 1.2,
              p["title"], size=64, color=B.WHITE, align=PP_ALIGN.CENTER, bold=True)
    if p.get("subtitle"):
        _add_text(s, 0.5, 3.65, 12.33, 0.8,
                  p["subtitle"], size=28, color=B.WHITE, align=PP_ALIGN.CENTER)
    if p.get("meta_line"):
        _add_text(s, 0.5, 5.3, 12.33, 0.4,
                  p["meta_line"], size=14, color=B.TEAL, align=PP_ALIGN.CENTER)
    if p.get("client_tag"):
        _add_text(s, 0.5, 5.75, 12.33, 0.4,
                  p["client_tag"], size=14, color=B.WHITE, align=PP_ALIGN.CENTER, bold=True)
    if p.get("run_date"):
        _add_text(s, 0.5, 6.15, 12.33, 0.4,
                  p["run_date"], size=12, color=B.WHITE, align=PP_ALIGN.CENTER)
    _add_text(s, 0.5, 7.0, 12.33, 0.4, "sumitup.",
              size=18, bold=True, color=B.WHITE, align=PP_ALIGN.CENTER)
    return page_num  # no increment on cover


def tpl_section_divider(prs, page_num, p):
    s = _new_navy_slide(prs)
    _add_text(s, 0, 3.2, 13.33, 1.3, p["title"],
              size=44, bold=True, color=B.WHITE, align=PP_ALIGN.CENTER)
    if p.get("subtitle"):
        _add_text(s, 0, 4.4, 13.33, 0.5, p["subtitle"],
                  size=18, italic=True, color=B.TEAL, align=PP_ALIGN.CENTER)
    return page_num  # no increment on divider


def tpl_executive_summary(prs, page_num, p):
    s = _new_content_slide(prs)
    page_num += 1
    _add_title(s, "Executive summary")
    if p.get("subtitle"):
        _add_subhead(s, p["subtitle"])
    frames = p["frames"]
    y0 = 1.8
    for i, f in enumerate(frames):
        y = y0 + i * 0.55
        _add_text(s, 0.5, y, 0.4, 0.5, str(f["num"]),
                  size=20, bold=True, color=B.TEAL)
        _add_text(s, 0.95, y, 4.8, 0.5, f["name"],
                  size=14, bold=True, color=B.NAVY)
        tag = f.get("tag", "")
        _add_text(s, 5.8, y + 0.05, 1.4, 0.4, tag,
                  size=10, bold=True, color=TAG_COLORS.get(tag, B.NAVY))
        _add_text(s, 7.3, y, 5.6, 0.5, f.get("description", ""),
                  size=11, color=B.NAVY)
    if p.get("source_line"):
        _add_source_citation(s, p["source_line"])
    _add_footer(s, page_num)
    return page_num


def tpl_outlet_landscape(prs, page_num, p):
    s = _new_content_slide(prs)
    page_num += 1
    _add_title(s, p.get("title", "Outlet landscape"))
    if p.get("subtitle"):
        _add_subhead(s, p["subtitle"])

    blocks = p["outlet_blocks"]
    color_map = {"navy": B.NAVY, "purple": B.PURPLE, "teal": B.TEAL, "gold": B.GOLD}
    col_w = 12.33 / max(len(blocks), 1)
    y0 = 2.0
    for i, blk in enumerate(blocks):
        x = 0.5 + i * col_w
        c = color_map.get(blk.get("color", "navy"), B.NAVY)
        _add_text(s, x, y0, col_w - 0.15, 0.5,
                  blk["position"].upper(),
                  size=12, bold=True, color=c, align=PP_ALIGN.LEFT)
        outlets_str = " · ".join(blk["outlets"])
        _add_text(s, x, y0 + 0.5, col_w - 0.15, 0.8, outlets_str,
                  size=14, bold=True, color=B.NAVY)
        _add_text(s, x, y0 + 1.6, col_w - 0.15, 2.8,
                  blk.get("signature_framing", ""),
                  size=11, color=B.NAVY)
    _add_footer(s, page_num)
    return page_num


def tpl_frame_argument(prs, page_num, p):
    s = _new_content_slide(prs)
    page_num += 1
    _add_title(s, f"Frame {p['frame_number']} · {p['frame_name']}")
    if p.get("tag"):
        _add_subhead(s, f"{p['tag']} · {p.get('one_line_summary', '')}")
    else:
        _add_subhead(s, p.get("one_line_summary", ""))

    # Two quote blocks side by side
    mq = p["main_quote"]
    sq = p.get("supporting_quote")
    if sq:
        _hebrew_quote_block(s, 0.5, 2.0, 6.0,
                            mq["hebrew"], mq["gloss"], mq["source"], h_heb=1.2)
        _hebrew_quote_block(s, 6.8, 2.0, 6.0,
                            sq["hebrew"], sq["gloss"], sq["source"], h_heb=1.2)
    else:
        _hebrew_quote_block(s, 0.5, 2.0, 12.33,
                            mq["hebrew"], mq["gloss"], mq["source"], h_heb=1.2)

    # Outlet presence
    y = 4.2
    op = p.get("outlet_presence", {})
    _add_text(s, 0.5, y, 12.33, 0.35, "Outlet presence:",
              size=12, bold=True, color=B.TEAL)
    lead = "LEAD: " + ", ".join(op.get("lead", []))
    secondary = "SECONDARY: " + ", ".join(op.get("secondary", []))
    _add_text(s, 0.5, y + 0.35, 12.33, 0.35, lead, size=11, color=B.NAVY)
    _add_text(s, 0.5, y + 0.7, 12.33, 0.35, secondary, size=11, color=B.NAVY)

    # Implication
    if p.get("implication"):
        _add_text(s, 0.5, 5.4, 12.33, 0.35, "Implication:",
                  size=12, bold=True, color=B.PURPLE)
        _add_text(s, 0.5, 5.75, 12.33, 0.8, p["implication"],
                  size=11, color=B.NAVY, italic=True)

    if p.get("sources"):
        src = "Examples: " + " · ".join(p["sources"][:3])
        _add_source_citation(s, src)
    _add_footer(s, page_num)
    return page_num


def tpl_frame_language(prs, page_num, p):
    s = _new_content_slide(prs)
    page_num += 1
    _add_title(s, f"Frame {p['frame_number']} · The Hebrew language")
    _add_subhead(s, f"Vocabulary carrying the '{p['frame_name']}' frame")

    phrases = p["phrases"]
    # 2-column grid
    col_w = 6.0
    row_h = 1.05
    for i, ph in enumerate(phrases[:10]):  # cap at 10
        col = i % 2
        row = i // 2
        x = 0.5 + col * (col_w + 0.3)
        y = 1.85 + row * row_h
        # Hebrew phrase (RTL)
        tb = _add_text(s, x, y, col_w, 0.38,
                       f'"{ph["hebrew"]}"',
                       size=13, bold=True, color=B.NAVY,
                       align=PP_ALIGN.RIGHT, rtl=True)
        # Gloss + valence
        val_color = {"positive": B.TEAL, "negative": B.PURPLE,
                     "contested": B.GOLD, "neutral": B.NAVY}.get(ph.get("valence", "neutral"), B.NAVY)
        _add_text(s, x, y + 0.36, col_w, 0.3,
                  f'{ph["gloss"]} · {ph.get("valence", "")}',
                  size=10, italic=True, color=val_color)
        # In-situ example
        if ph.get("in_situ"):
            outlet = ph.get("in_situ_outlet", "")
            _add_text(s, x, y + 0.64, col_w, 0.38,
                      f'({outlet}) {ph["in_situ"]}',
                      size=9, color=B.NAVY, italic=True,
                      align=PP_ALIGN.RIGHT, rtl=True)

    if p.get("note"):
        _add_text(s, 0.5, 6.45, 12.33, 0.3,
                  p["note"], size=11, italic=True, color=B.PURPLE)
    _add_footer(s, page_num)
    return page_num


def tpl_frame_positions(prs, page_num, p):
    s = _new_content_slide(prs)
    page_num += 1
    _add_title(s, f"Frame {p['frame_number']} · Ideological positions")
    _add_subhead(s, f"How different outlets frame '{p['frame_name']}'")

    positions = p["positions"]
    n = len(positions)
    col_w = (12.33 - 0.3 * (n - 1)) / n
    y_title = 1.95
    y_quote = 3.2

    for i, pos in enumerate(positions[:3]):
        x = 0.5 + i * (col_w + 0.3)
        color = [B.NAVY, B.PURPLE, B.TEAL][i]
        # Position name
        _add_text(s, x, y_title, col_w, 0.4,
                  pos["name"].upper(), size=13, bold=True, color=color)
        # Outlets
        _add_text(s, x, y_title + 0.38, col_w, 0.3,
                  pos.get("outlets", ""), size=10, italic=True, color=B.NAVY)
        # Framing valence
        _add_text(s, x, y_title + 0.72, col_w, 0.35,
                  pos.get("framing_valence", ""),
                  size=11, bold=True, color=color, italic=True)
        # Quote block
        q = pos.get("quote")
        if q:
            _hebrew_quote_block(s, x, y_quote, col_w,
                                q.get("hebrew", ""),
                                q.get("gloss", ""),
                                f'{q.get("speaker", "")}, {q.get("outlet_date", "")}',
                                h_heb=1.3)

    # Absences note + convergence note
    y_foot = 5.8
    if p.get("absences"):
        _add_text(s, 0.5, y_foot, 12.33, 0.35,
                  "Absent from this frame:",
                  size=11, bold=True, color=B.GRAY)
        _add_text(s, 0.5, y_foot + 0.3, 12.33, 0.35,
                  p["absences"], size=10, italic=True, color=B.GRAY)
    if p.get("convergence_note"):
        _add_text(s, 0.5, 6.35, 12.33, 0.35,
                  p["convergence_note"],
                  size=11, italic=True, color=B.PURPLE)
    _add_footer(s, page_num)
    return page_num


def tpl_vocabulary_cluster(prs, page_num, p):
    s = _new_content_slide(prs)
    page_num += 1
    valence_color = {"positive": B.TEAL, "negative": B.PURPLE,
                     "contested": B.GOLD, "neutral": B.NAVY}.get(
        p.get("cluster_valence", "neutral"), B.NAVY)
    _add_title(s, f"Vocabulary · {p['cluster_name']}")
    if p.get("cluster_description"):
        _add_subhead(s, p["cluster_description"], color=valence_color)

    # 2-column phrase grid
    phrases = p["phrases"]
    col_w = 6.0
    row_h = 0.5
    max_rows_per_col = 8  # 16 phrases max
    for i, ph in enumerate(phrases[:16]):
        col = i % 2
        row = i // 2
        x = 0.5 + col * (col_w + 0.3)
        y = 1.9 + row * row_h
        # Hebrew (RTL) + gloss on one line
        tb = _add_text(s, x, y, col_w, 0.32,
                       f'"{ph["hebrew"]}"',
                       size=13, bold=True, color=B.NAVY,
                       align=PP_ALIGN.RIGHT, rtl=True)
        _add_text(s, x, y + 0.28, col_w, 0.22,
                  f'{ph["gloss"]} · ({ph.get("outlets", "")})',
                  size=9, italic=True, color=valence_color)

    # Signature quote at bottom
    sq = p.get("signature_quote")
    if sq:
        y_sig = 6.15
        _add_text(s, 0.5, y_sig, 6.0, 0.3,
                  "Signature quote:", size=10, bold=True, color=B.PURPLE)
        _add_text(s, 0.5, y_sig + 0.28, 12.33, 0.3,
                  f'"{sq.get("hebrew", "")}" — {sq.get("source", "")}',
                  size=10, color=B.NAVY, italic=True,
                  align=PP_ALIGN.RIGHT, rtl=True)
    _add_footer(s, page_num)
    return page_num


def tpl_temporal_phase(prs, page_num, p):
    s = _new_content_slide(prs)
    page_num += 1
    _add_title(s, f"Phase {p['phase_num']} · {p['phase_title']}")
    _add_subhead(s, p.get("date_range", ""))

    _add_text(s, 0.5, 1.7, 12.33, 0.6,
              p.get("phase_description", ""),
              size=12, color=B.NAVY, italic=True)

    examples = p.get("article_examples", [])[:3]
    col_w = 4.0
    for i, ex in enumerate(examples):
        x = 0.5 + i * (col_w + 0.15)
        y = 2.6
        # Outlet + date
        _add_text(s, x, y, col_w, 0.3,
                  f'{ex["outlet"]} · {ex["date"]}',
                  size=10, bold=True, color=B.TEAL)
        # Title excerpt
        _add_text(s, x, y + 0.3, col_w, 0.7,
                  ex.get("title", ""),
                  size=11, bold=True, color=B.NAVY)
        # Hebrew quote
        tb = _add_text(s, x, y + 1.05, col_w, 1.8,
                       f'"{ex.get("hebrew_quote", "")}"',
                       size=10, color=B.NAVY,
                       align=PP_ALIGN.RIGHT, rtl=True)
        # Gloss
        _add_text(s, x, y + 2.9, col_w, 0.8,
                  f'"{ex.get("english_gloss", "")}"',
                  size=9, italic=True, color=B.NAVY)

    if p.get("signature_finding"):
        _add_text(s, 0.5, 6.2, 12.33, 0.4,
                  p["signature_finding"],
                  size=11, italic=True, color=B.PURPLE)
    _add_footer(s, page_num)
    return page_num


def tpl_battleground(prs, page_num, p):
    s = _new_content_slide(prs)
    page_num += 1
    _add_title(s, f"Battleground · {p['battleground_name']}")
    _add_subhead(s, p.get("context", ""))

    has_third = p.get("third") is not None
    n = 3 if has_third else 2
    col_w = (12.33 - 0.3 * (n - 1)) / n

    sides = [p["left"], p["right"]]
    colors = [B.NAVY, B.PURPLE]
    if has_third:
        sides.append(p["third"])
        colors.append(B.TEAL)

    for i, side in enumerate(sides):
        x = 0.5 + i * (col_w + 0.3)
        c = colors[i]
        _add_text(s, x, 1.9, col_w, 0.4,
                  side["position"].upper(), size=13, bold=True, color=c)
        _add_text(s, x, 2.3, col_w, 0.3,
                  f'{side.get("outlet", "")} · {side.get("speaker", "")}',
                  size=10, italic=True, color=B.NAVY)
        # Hebrew quote (large, main content)
        tb = _add_text(s, x, 2.75, col_w, 2.3,
                       f'"{side["hebrew_quote"]}"',
                       size=12, bold=True, color=B.NAVY,
                       align=PP_ALIGN.RIGHT, rtl=True)
        # Gloss
        _add_text(s, x, 5.1, col_w, 0.9,
                  f'"{side.get("english_gloss", "")}"',
                  size=10, italic=True, color=B.NAVY)
        # Rhetorical move
        _add_text(s, x, 6.0, col_w, 0.3,
                  side.get("rhetorical_move", ""),
                  size=10, bold=True, color=c, italic=True)

    if p.get("axis_of_contestation"):
        _add_text(s, 0.5, 6.5, 12.33, 0.3,
                  "The axis: " + p["axis_of_contestation"],
                  size=10, italic=True, color=B.PURPLE)
    _add_footer(s, page_num)
    return page_num


def tpl_experts_grid(prs, page_num, p):
    s = _new_content_slide(prs)
    page_num += 1
    _add_title(s, "Who speaks · cast of characters")
    if p.get("subtitle"):
        _add_subhead(s, p["subtitle"])

    experts = p.get("experts", [])[:9]  # 3×3 grid max per slide
    col_w = 4.0
    row_h = 1.7
    y0 = 1.85
    for i, ex in enumerate(experts):
        col = i % 3
        row = i // 3
        x = 0.5 + col * (col_w + 0.15)
        y = y0 + row * row_h
        _add_text(s, x, y, col_w, 0.3,
                  ex["name"], size=12, bold=True, color=B.NAVY)
        _add_text(s, x, y + 0.3, col_w, 0.3,
                  ex.get("credential", ""), size=10, italic=True, color=B.TEAL)
        _add_text(s, x, y + 0.6, col_w, 0.3,
                  ex.get("institution", ""), size=10, color=B.NAVY)
        _add_text(s, x, y + 0.9, col_w, 0.3,
                  f'{ex.get("article_count", "")} articles · {ex.get("representational_note", "")}',
                  size=9, italic=True, color=B.NAVY)
        if ex.get("representative_hebrew_quote"):
            _add_text(s, x, y + 1.2, col_w, 0.45,
                      f'"{ex["representative_hebrew_quote"]}"',
                      size=9, color=B.NAVY, italic=True,
                      align=PP_ALIGN.RIGHT, rtl=True)

    if p.get("scope_note"):
        _add_text(s, 0.5, 6.4, 12.33, 0.3,
                  p["scope_note"], size=10, italic=True, color=B.NAVY)
    _add_footer(s, page_num)
    return page_num


def tpl_quote_vs_avoid(prs, page_num, p):
    s = _new_content_slide(prs)
    page_num += 1
    _add_title(s, "Voices to quote vs. voices to avoid")
    if p.get("subtitle"):
        _add_subhead(s, p["subtitle"])

    # Left: quote side (teal)
    q = p["quote_side"]
    _add_text(s, 0.5, 1.9, 6.0, 0.4,
              "\u2713  " + q.get("label", "QUOTE · ALIGN WITH"),
              size=14, bold=True, color=B.TEAL)
    for i, entry in enumerate(q.get("entries", [])[:6]):
        y = 2.4 + i * 0.7
        _add_text(s, 0.5, y, 6.0, 0.3,
                  "\u2022 " + entry["name"], size=12, bold=True, color=B.NAVY)
        _add_text(s, 0.75, y + 0.3, 5.75, 0.35,
                  entry["reason"], size=10, italic=True, color=B.NAVY)

    # Right: avoid side (purple)
    a = p["avoid_side"]
    _add_text(s, 6.8, 1.9, 6.0, 0.4,
              "\u2717  " + a.get("label", "AVOID"),
              size=14, bold=True, color=B.PURPLE)
    for i, entry in enumerate(a.get("entries", [])[:6]):
        y = 2.4 + i * 0.7
        _add_text(s, 6.8, y, 6.0, 0.3,
                  "\u2022 " + entry["name"], size=12, bold=True, color=B.NAVY)
        _add_text(s, 7.05, y + 0.3, 5.75, 0.35,
                  entry["reason"], size=10, italic=True, color=B.NAVY)

    _add_footer(s, page_num)
    return page_num


def tpl_absence_detail(prs, page_num, p):
    s = _new_content_slide(prs)
    page_num += 1
    conf = p.get("confidence", "").upper()
    conf_color = {"HIGH": B.TEAL, "MEDIUM": B.PURPLE, "LOW": B.GRAY}.get(conf, B.NAVY)
    _add_title(s, f"Absent · {p['absence_name']}")
    _add_subhead(s, f"{p.get('absence_category', '')} · {conf} confidence", color=conf_color)

    y = 1.9
    rows = [
        ("What's missing", p.get("what_is_missing", ""), B.NAVY),
        ("Why expected", p.get("why_expected", ""), B.NAVY),
        ("Evidence", p.get("evidence", ""), B.NAVY),
    ]
    for label, content, c in rows:
        _add_text(s, 0.5, y, 2.5, 0.4, label + ":",
                  size=12, bold=True, color=B.TEAL)
        _add_text(s, 3.2, y, 9.6, 1.1, content,
                  size=12, color=c)
        y += 1.25

    if p.get("strategic_implication"):
        _add_text(s, 0.5, 5.9, 12.33, 0.4,
                  "Strategic implication:",
                  size=12, bold=True, color=B.PURPLE)
        _add_text(s, 0.5, 6.25, 12.33, 0.5,
                  p["strategic_implication"],
                  size=11, italic=True, color=B.NAVY)
    _add_footer(s, page_num)
    return page_num


def tpl_client_white_space(prs, page_num, p):
    s = _new_content_slide(prs)
    page_num += 1
    _add_title(s, p.get("title", "White space"))
    if p.get("subtitle"):
        _add_subhead(s, p["subtitle"])
    territories = p.get("territories", [])[:3]
    y = 1.85
    for t in territories:
        c = [B.TEAL, B.PURPLE, B.NAVY][int(t.get("num", 1)) - 1] if t.get("num") else B.TEAL
        _add_text(s, 0.5, y, 0.5, 0.5, str(t.get("num", "")),
                  size=28, bold=True, color=c)
        _add_text(s, 1.15, y, 11.7, 0.4,
                  t.get("name", ""), size=14, bold=True, color=B.NAVY)
        _add_text(s, 1.15, y + 0.42, 11.7, 0.8,
                  t.get("description", ""), size=11, color=B.NAVY)
        if t.get("first_move"):
            _add_text(s, 1.15, y + 1.2, 11.7, 0.35,
                      "First move: " + t["first_move"],
                      size=10, italic=True, color=c)
        y += 1.65
    _add_footer(s, page_num)
    return page_num


def tpl_client_plays(prs, page_num, p):
    s = _new_content_slide(prs)
    page_num += 1
    _add_title(s, p.get("title", "Positioning plays"))
    plays = p.get("plays", [])[:5]
    y = 1.8
    for i, play in enumerate(plays):
        _add_text(s, 0.5, y, 0.5, 0.5, str(i + 1),
                  size=22, bold=True, color=B.TEAL)
        _add_text(s, 1.1, y, 11.7, 0.35,
                  play.get("frame", ""), size=12, bold=True, color=B.NAVY)
        _add_text(s, 1.1, y + 0.35, 11.7, 0.35,
                  "\u2192 " + play.get("play", ""),
                  size=11, bold=True, color=B.PURPLE)
        vocab = play.get("vocabulary", [])
        if vocab:
            vocab_str = " · ".join([f'"{v}"' for v in vocab])
            _add_text(s, 1.1, y + 0.68, 11.7, 0.35,
                      "Vocab: " + vocab_str, size=10, italic=True, color=B.NAVY,
                      align=PP_ALIGN.RIGHT, rtl=True)
        if play.get("voice"):
            _add_text(s, 1.1, y + 1.0, 11.7, 0.3,
                      "Voice: " + play["voice"], size=10, italic=True, color=B.TEAL)
        y += 1.05
    _add_footer(s, page_num)
    return page_num


def tpl_client_avoids(prs, page_num, p):
    s = _new_content_slide(prs)
    page_num += 1
    _add_title(s, p.get("title", "What to avoid"))
    avoids = p.get("avoids", [])[:5]
    y = 1.85
    for av in avoids:
        _add_text(s, 0.5, y, 0.5, 0.5, "\u2717",
                  size=22, bold=True, color=B.PURPLE)
        _add_text(s, 1.0, y, 11.83, 0.35,
                  av.get("name", ""), size=13, bold=True, color=B.NAVY)
        _add_text(s, 1.0, y + 0.35, 11.83, 0.7,
                  av.get("reason", ""), size=11, color=B.NAVY)
        if av.get("evidence_anchor"):
            _add_text(s, 1.0, y + 1.05, 11.83, 0.35,
                      "Evidence: " + av["evidence_anchor"],
                      size=10, italic=True, color=B.TEAL)
        y += 1.5
    _add_footer(s, page_num)
    return page_num


def tpl_client_taglines(prs, page_num, p):
    s = _new_content_slide(prs)
    page_num += 1
    _add_title(s, p.get("title", "Hebrew taglines worth testing"))
    taglines = p.get("taglines", [])[:3]
    y = 2.0
    for i, t in enumerate(taglines):
        c = [B.TEAL, B.PURPLE, B.NAVY][i]
        _add_text(s, 0.5, y, 0.5, 0.5, str(i + 1),
                  size=24, bold=True, color=c)
        # Hebrew RTL
        tb = _add_text(s, 1.1, y, 11.7, 0.5,
                       f'"{t.get("hebrew", "")}"',
                       size=20, bold=True, color=B.NAVY,
                       align=PP_ALIGN.RIGHT, rtl=True)
        _add_text(s, 1.1, y + 0.5, 11.7, 0.3,
                  f'"{t.get("gloss", "")}"',
                  size=11, italic=True, color=B.NAVY)
        _add_text(s, 1.1, y + 0.8, 11.7, 0.7,
                  t.get("why_it_fits", ""), size=10, color=c, italic=True)
        y += 1.55
    _add_footer(s, page_num)
    return page_num


def tpl_methodology(prs, page_num, p):
    s = _new_content_slide(prs)
    page_num += 1
    _add_title(s, "Methodology & limitations")
    if p.get("subtitle"):
        _add_subhead(s, p["subtitle"])

    # Left: what we did
    _add_text(s, 0.5, 1.9, 6.0, 0.35, "What we did",
              size=12, bold=True, color=B.TEAL)
    _add_text(s, 0.5, 2.3, 6.0, 3.5,
              ["\u2022 " + item for item in p.get("what_we_did", [])],
              size=10, color=B.NAVY)

    # Right: verifier contests
    _add_text(s, 6.8, 1.9, 5.83, 0.35, "Verifier-contested claims (resolved)",
              size=12, bold=True, color=B.PURPLE)
    contests = p.get("verifier_contests", [])
    y = 2.3
    for c in contests[:5]:
        _add_text(s, 6.8, y, 5.83, 0.5, "\u2022 " + c.get("claim", ""),
                  size=10, bold=True, color=B.NAVY)
        _add_text(s, 7.0, y + 0.5, 5.63, 0.5, c.get("resolution", ""),
                  size=9, italic=True, color=B.NAVY)
        y += 1.05

    # Bottom: scope limits
    _add_text(s, 0.5, 5.95, 12.33, 0.35, "What this scan is NOT",
              size=12, bold=True, color=B.NAVY)
    _add_text(s, 0.5, 6.3, 12.33, 0.7,
              ["\u2022 " + item for item in p.get("scope_limits", [])],
              size=10, color=B.NAVY)
    _add_footer(s, page_num)
    return page_num


def tpl_closing(prs, page_num, p):
    s = _new_navy_slide(prs)
    _add_text(s, 0, 2.8, 13.33, 1.2, p.get("title", "Thank you."),
              size=64, bold=True, color=B.WHITE, align=PP_ALIGN.CENTER)
    if p.get("subtitle"):
        _add_text(s, 0, 4.2, 13.33, 0.5, p["subtitle"],
                  size=18, color=B.WHITE, align=PP_ALIGN.CENTER, italic=True)
    if p.get("client_line"):
        _add_text(s, 0, 4.8, 13.33, 0.4, p["client_line"],
                  size=14, color=B.TEAL, align=PP_ALIGN.CENTER)
    _add_text(s, 0, 6.7, 13.33, 0.4, "sumitup.",
              size=24, bold=True, color=B.WHITE, align=PP_ALIGN.CENTER)
    return page_num  # no increment


def tpl_appendix_cover(prs, page_num, p):
    s = _new_navy_slide(prs)
    _add_text(s, 0, 3.0, 13.33, 1.0, p.get("title", "Appendix"),
              size=54, bold=True, color=B.WHITE, align=PP_ALIGN.CENTER)
    if p.get("subtitle"):
        _add_text(s, 0, 4.0, 13.33, 0.5, p["subtitle"],
                  size=18, color=B.WHITE, align=PP_ALIGN.CENTER, italic=True)
    return page_num  # no increment


def tpl_article_index_page(prs, page_num, p):
    s = _new_content_slide(prs)
    page_num += 1
    _add_title(s, f"Articles · {p.get('frame', 'index')}")
    count = p.get("article_count_on_this_frame", "")
    if count:
        _add_subhead(s, f"{count} articles — primary or secondary framing")

    articles = p.get("articles", [])[:18]  # ~18 per page
    # 2-column layout
    col_w = 6.0
    row_h = 0.3
    for i, a in enumerate(articles):
        col = i % 2
        row = i // 2
        x = 0.5 + col * (col_w + 0.3)
        y = 1.9 + row * row_h
        line = f'{a.get("outlet", "")} · {a.get("date", "")} · {a.get("title_excerpt", "")[:60]}'
        tag = a.get("primary_or_secondary", "")
        color = B.NAVY if tag == "primary" else B.GRAY
        _add_text(s, x, y, col_w, row_h, line, size=9, color=color)

    _add_footer(s, page_num)
    return page_num


# ============================================================
# TEMPLATE REGISTRY
# ============================================================
TEMPLATES = {
    "cover":              tpl_cover,
    "section_divider":    tpl_section_divider,
    "executive_summary":  tpl_executive_summary,
    "outlet_landscape":   tpl_outlet_landscape,
    "frame_argument":     tpl_frame_argument,
    "frame_language":     tpl_frame_language,
    "frame_positions":    tpl_frame_positions,
    "vocabulary_cluster": tpl_vocabulary_cluster,
    "temporal_phase":     tpl_temporal_phase,
    "battleground":       tpl_battleground,
    "experts_grid":       tpl_experts_grid,
    "quote_vs_avoid":     tpl_quote_vs_avoid,
    "absence_detail":     tpl_absence_detail,
    "client_white_space": tpl_client_white_space,
    "client_plays":       tpl_client_plays,
    "client_avoids":      tpl_client_avoids,
    "client_taglines":    tpl_client_taglines,
    "methodology":        tpl_methodology,
    "closing":            tpl_closing,
    "appendix_cover":     tpl_appendix_cover,
    "article_index_page": tpl_article_index_page,
}


# ============================================================
# PRESENTATION BOOTSTRAP
# ============================================================

def new_presentation():
    """Create a new 16:9 presentation, cream-default."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def dispatch_slide(prs, page_num, slide_entry):
    """Render one slide from a deck_plan.yaml entry. Returns new page_num."""
    template_name = slide_entry["template"]
    params = slide_entry.get("params", {})
    fn = TEMPLATES.get(template_name)
    if fn is None:
        raise ValueError(f"Unknown template: {template_name!r}. "
                         f"Available: {sorted(TEMPLATES.keys())}")
    return fn(prs, page_num, params)
